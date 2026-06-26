"""
PopCart PH — DARK narrative pitch deck (LAMBO 2026, PDF submission).
Structured to the case's A-E output format; every fact cited to its exhibit;
written to be read (no presenter). Run: .venv/bin/python build_deck.py
"""
import struct
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0x0B,0x0E,0x11); TITLE=RGBColor(0xFF,0xFF,0xFF); LIME=RGBColor(0xC3,0xF5,0x3C)
ORANGE=RGBColor(0xF3,0x9C,0x12); INK=RGBColor(0xD6,0xDB,0xE0); SOFT=RGBColor(0xEC,0xEF,0xF2)
MUTE=RGBColor(0x8A,0x92,0x9B); SLATE=RGBColor(0x16,0x21,0x2E); SLATE2=RGBColor(0x1E,0x2A,0x3A)
BORDER=RGBColor(0x2A,0x38,0x48); GREEN=RGBColor(0x2E,0xCC,0x71); RED=RGBColor(0xE7,0x4C,0x3C)
SKY=RGBColor(0x5B,0x9B,0xD5); GOLD=RGBColor(0xF1,0xC4,0x0F)
SECT={"A":"A · STRATEGIC RECOMMENDATION","B":"B · WHY IT MAKES SENSE",
      "C":"C · HOW IT WILL WORK","D":"D · FINANCIAL LOGIC","E":"E · ROADMAP & RISK"}
SECTC={"A":LIME,"B":SKY,"C":GOLD,"D":GREEN,"E":ORANGE}

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; NUM=[0]
def png(p):
    with open(p,"rb") as f: h=f.read(24)
    return struct.unpack(">II",h[16:24])
def slide():
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=BG; r.line.fill.background(); r.shadow.inherit=False
    return s
def rect(s,x,y,w,h,fc=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fc is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fc
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def text(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for i,(t,sz,b,c,*rest) in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if rest: p.space_after=Pt(rest[0])
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c
        r.font.italic=italic; r.font.name="Arial"
    return tb
def rich(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor; p=tf.paragraphs[0]; p.alignment=align
    for t,sz,b,c in runs:
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c
        r.font.italic=italic; r.font.name="Arial"
    return tb
def title2(s,segs,sect=None,y=0.64,size=31):
    if sect: text(s,0.58,0.3,11.5,0.3,[(SECT[sect],11,True,SECTC[sect])])
    tb=s.shapes.add_textbox(Inches(0.55),Inches(y),Inches(12.2),Inches(0.9)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    for t,c in segs:
        r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.bold=True; r.font.color.rgb=c; r.font.name="Arial"
    rect(s,0.57,y+0.82,2.1,0.045,LIME)
def voice(s,txt,y=1.62):
    rect(s,0.57,y+0.03,0.06,0.33,LIME); text(s,0.78,y,11.9,0.5,[(txt,15,False,SOFT)],italic=True)
def bullets(s,x,y,w,h,items,sz=13.5,gap=10,color=INK):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        a=p.add_run(); a.text="●  "; a.font.size=Pt(sz-4); a.font.color.rgb=LIME; a.font.bold=True; a.font.name="Arial"
        b=p.add_run(); b.text=it; b.font.size=Pt(sz); b.font.color.rgb=color; b.font.name="Arial"
def img(s,path,x,y,w,h,frame=False,caption=None):
    pw,ph=png(path); a=pw/ph
    if w/h>a: nh=h; nw=h*a
    else: nw=w; nh=w/a
    px=x+(w-nw)/2; py=y+(h-nh)/2
    if frame: rect(s,px-0.12,py-0.12,nw+0.24,nh+0.24,SLATE,line=BORDER,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    s.shapes.add_picture(path,Inches(px),Inches(py),Inches(nw),Inches(nh))
    if caption: text(s,x,py+nh+0.05,w,0.3,[(caption,10,False,MUTE)],align=PP_ALIGN.CENTER)
def source(s,txt,y=6.62):
    text(s,0.57,y,12.2,0.3,[("Source — "+txt,10,False,MUTE)],italic=True)
def footer(s):
    NUM[0]+=1
    text(s,0.57,7.06,9,0.3,[("PopCart PH   ·   LAMBO 2026   ·   [School] · [Team]",9,False,MUTE)])
    text(s,12.0,7.06,0.8,0.3,[(f"{NUM[0]:02d}",9,True,MUTE)],align=PP_ALIGN.RIGHT)
def notes(s,t): s.notes_slide.notes_text_frame.text=t
def statcard(s,x,y,w,h,big,bigc,lab,sub):
    rect(s,x,y,w,h,SLATE,line=BORDER,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s,x,y,w,0.07,bigc)
    text(s,x+0.25,y+0.2,w-0.5,0.9,[(big,32,True,bigc)])
    text(s,x+0.27,y+1.0,w-0.5,0.4,[(lab,12.5,True,SOFT)])
    text(s,x+0.27,y+1.38,w-0.5,h-1.42,[(sub,11,False,MUTE)])

# ============ COVER ============
s=slide()
text(s,0.8,1.4,11.7,0.4,[("POPCART PH   ·   LAMBO 2026 QUALIFYING CASE",13,True,MUTE)])
rich(s,0.8,2.0,12.4,1.0,[("FIX THE ",44,True,TITLE),("ENGINE",44,True,LIME)])
text(s,0.8,3.02,12.4,1.0,[("BEFORE YOU GROW THE BODY",44,True,TITLE)])
rect(s,0.83,4.42,2.4,0.05,ORANGE)
text(s,0.8,4.68,11.4,1.1,[("PopCart is profitable, viral, and growing fast — ₱120M in sales at a 9% margin "
    "(Exhibit A) — and quietly losing two-thirds of its profit to a leak no one is watching.",18,False,SOFT)])
text(s,0.8,6.0,11.4,0.5,[("Here's where the money goes — and the cheapest way to stop it.",15,True,ORANGE)],italic=True)
text(s,0.8,6.82,11.7,0.4,[("[School Name]   ·   [Team Name]",12,False,MUTE)])
notes(s,"Written deck, no presenter. The cover states the paradox (healthy + bleeding) and promises a payoff.")

# ============ S1 · A — BOTTOM LINE ============
s=slide(); title2(s,[("THE ",TITLE),("BOTTOM LINE",LIME)],sect="A")
voice(s,"If you read one slide, read this one — the entire case in five answers.")
rect(s,0.57,2.25,12.2,0.72,SLATE,line=LIME,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rich(s,0.85,2.33,11.7,0.6,[("OUR CALL:  ",14,True,LIME),("Fix operations first. Expand only in sequence — "
    "never all at once.",14,True,SOFT)],anchor=MSO_ANCHOR.MIDDLE)
qa=[("Q1 · Biggest problem","Stockouts — ₱7.35M/yr in lost gross profit, not shipping  (Exhibit C)"),
    ("Q2 · Where to expand","Nowhere yet — fix ops online; a Cebu 3PL is gated to Month 7"),
    ("Q3 · The tool","A simple dashboard that gives earlier reorder warnings  (improves Exhibit D)"),
    ("Q4 · The impact","Stockout days 11 → ~6/mo; inventory updates 2×/day → every 2h  (Exhibit C/D)"),
    ("Q5 · Worth it?","Adds ₱1.0–1.6M recurring benefit/yr; ₱1.4M setup; ~13–20 mo payback")]
y=3.25
for tag,ans in qa:
    rich(s,0.7,y,12.2,0.5,[(tag+"     ",13,True,LIME),(ans,13,False,INK)]); y+=0.52
rect(s,0.57,5.95,12.2,0.6,SLATE2,line=BORDER,lw=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rich(s,0.6,6.03,12.1,0.45,[("₱7.35M",15,True,RED),(" gross-profit leak  ·  ",13,False,MUTE),("+₱0.20M",15,True,GREEN),
    (" cautious recurring floor  ·  ",13,False,MUTE),("13–20 mo",15,True,SKY),(" payback",13,False,MUTE)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
footer(s); notes(s,"Answer first (BLUF). Everything after this is evidence, in the case's A-E order.")

# ============ S2 · B — FIVE PROBLEMS ============
s=slide(); title2(s,[("FIVE PROBLEMS, ",TITLE),("ONE QUESTION",LIME)],sect="B")
voice(s,"Before choosing a fix, we asked a sharper question: which problem is quietly the most expensive?")
probs=[("Shipping to Visayas / Mindanao is costly","₱155 / ₱185 per order","Exhibit B"),
       ("Delivery is slow","5–7 days (Cebu), 7–10 days (Davao)","Exhibit B"),
       ("Products arrive damaged","4.2% / 5.0% return rate","Exhibit B"),
       ("Stockouts on viral best-sellers","out of stock up to 11 days/mo","Exhibit C"),
       ("90% of sales on one channel","TikTok Shop dependence","Exhibit D")]
y=2.35
for i,(p,d,ex) in enumerate(probs):
    rect(s,0.57,y,12.2,0.6,SLATE if i==3 else BG,line=BORDER if i==3 else None,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,0.8,y+0.05,0.5,0.5,[(str(i+1),16,True,LIME if i==3 else MUTE)],anchor=MSO_ANCHOR.MIDDLE)
    text(s,1.35,y+0.05,7.2,0.5,[(p,13.5,i==3,SOFT if i==3 else INK)],anchor=MSO_ANCHOR.MIDDLE)
    text(s,8.4,y+0.05,3.0,0.5,[(d,12,False,INK)],anchor=MSO_ANCHOR.MIDDLE)
    text(s,11.5,y+0.05,1.2,0.5,[(ex,10.5,True,SECTC["B"])],anchor=MSO_ANCHOR.MIDDLE)
    y+=0.66
text(s,0.57,6.35,12.2,0.5,[("Shipping is the loudest cost. Stockouts are the quietest quantified opportunity — "
    "and the cheapest first move to test.",12.5,True,SOFT)],align=PP_ALIGN.CENTER,italic=True)
footer(s); notes(s,"We frame the 'why' section as a hunt for the real culprit, and pre-highlight stockouts.")

# ============ S3 · B — FOLLOW THE MONEY ============
s=slide(); title2(s,[("FOLLOW THE ",TITLE),("MONEY",LIME)],sect="B")
voice(s,"Read Exhibit A with us — where does PopCart's ₱120M actually go?")
rows=[("Revenue","₱120.0M","100%",SOFT),("− Cost of goods sold","(₱60.0M)","50%",INK),
      ("− Platform fees & marketing","(₱12.0M)","10%",INK),
      ("− Logistics & shipping subsidies","(₱24.0M)","20%",ORANGE),
      ("− Warehouse OpEx + Admin","(₱13.2M)","11%",INK),("= Net profit","₱10.8M","9%",GREEN)]
y=2.4
for lab,amt,pct,c in rows:
    hl = c==ORANGE
    if hl: rect(s,0.57,y-0.03,7.6,0.5,SLATE,line=ORANGE,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,0.75,y,4.6,0.45,[(lab,13.5,c in(SOFT,GREEN,ORANGE),c)],anchor=MSO_ANCHOR.MIDDLE)
    text(s,5.3,y,1.7,0.45,[(amt,13.5,True,c)],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.RIGHT)
    text(s,7.1,y,1.0,0.45,[(pct,12,False,MUTE)],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.RIGHT)
    y+=0.55
rect(s,8.7,2.5,4.0,3.0,SLATE,line=BORDER,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,8.95,2.75,3.5,0.5,[("THE OBVIOUS INSTINCT",12.5,True,ORANGE)])
text(s,8.95,3.3,3.5,2.0,[("Shipping is 20% of revenue — the biggest controllable line.",13,False,SOFT,8),
    ("So cut it: build a regional hub.",13,True,SOFT,8),("It feels right.",13,False,INK,8),
    ("Watch what happens when we test it.  →",12.5,False,ORANGE)],)
source(s,"Exhibit A — Simplified Financial Snapshot, FY2025")
footer(s); notes(s,"We let the reader reach the 'build a hub' conclusion — then dismantle it on S5. The ponder setup.")

# ============ S4 · B — THE LEAK ============
s=slide(); title2(s,[("THE LEAK YOU ",TITLE),("CAN'T SEE",LIME)],sect="B")
voice(s,"Now read Exhibit C — and do this one calculation with us.")
bullets(s,0.57,2.35,6.1,3.4,[
 "Freeze-dried candy sells 300 units/day — but is out of stock 11 days a month  (Exhibit C).",
 "300 × 11 × 12 = 39,600 units never sold → × ₱180 = ₱7.13M → × 50% margin (Exhibit A) = ₱3.56M lost on ONE product.",
 "Across all fast-movers: ₱7.35M/yr in lost gross profit = 68% of net profit (Exhibit A).",
 "Exhibit D shows the control gap: stock updates happen by hand twice a day, while ordering relies on visual checks and gut feel."],sz=13.5,gap=13)
img(s,"assets/01_leak_bucket.png",6.7,2.15,6.1,4.0)
text(s,0.57,6.0,6.1,0.5,[("Shipping is the biggest cost. Stockouts are the biggest quantified recoverable opportunity.",
    12.5,True,LIME)],italic=True)
source(s,"Exhibit C (velocity) · Exhibit A (margin) · Exhibit D (manual ops)")
footer(s); notes(s,"The aha — the analytical and emotional center. The reader does the math and feels the loss.")

# ============ S5 · B — THE HUB TRAP ============
s=slide(); title2(s,[("WHY THE OBVIOUS MOVE ",TITLE),("BACKFIRES",LIME)],sect="B")
voice(s,"Remember the 'build a hub' instinct? Let's put real numbers on it.")
bullets(s,0.57,2.4,6.0,3.4,[
 "Central Visayas: 48,100 orders/yr (Assumption B) at ₱155/order (Exhibit B).",
 "A hub cuts shipping ~₱50/order (Assumption D) → 48,100 × ₱50 = ₱2.4M saved.",
 "Base damage savings (~₱0.29M) + 15% growth (~₱2.34M GP) lift total benefit to only ~₱5.04M.",
 "Against ₱9M annual operating cost, the base case loses ~₱3.96M.",
 "Break-even needs roughly 40% regional growth — above the case's 15–30% planning range."],sz=12.3,gap=9)
text(s,7.0,2.45,5.6,0.4,[("WHAT A HUB DOES TO THE MATH",12,True,MUTE)])
rect(s,7.0,3.05,2.2,0.5,GREEN,shape=MSO_SHAPE.ROUNDED_RECTANGLE); text(s,9.35,3.08,3.0,0.45,[("₱5.04M total benefit",12.5,True,GREEN)],anchor=MSO_ANCHOR.MIDDLE)
rect(s,7.0,3.95,4.0,0.5,RED,shape=MSO_SHAPE.ROUNDED_RECTANGLE); text(s,7.05,3.98,4.0,0.45,[("₱9M to run /yr",12.5,True,TITLE)],anchor=MSO_ANCHOR.MIDDLE)
rect(s,7.0,5.0,5.6,0.95,SLATE,line=ORANGE,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,7.2,5.12,5.3,0.75,[("Base case still loses ~₱3.96M / year",15,True,ORANGE)],anchor=MSO_ANCHOR.MIDDLE)
source(s,"Exhibit B · Assumptions B/C/D/E · Exhibit A")
footer(s); notes(s,"Payoff of S3: the reader proposed the hub; the data talks them out of it. Trust is earned.")

# ============ S6 · B — OPTIONS ============
s=slide(); title2(s,[("THE OPTIONS, ",TITLE),("SIDE BY SIDE",LIME)],sect="B")
voice(s,"We modeled each option's range — because on a 9% margin, the downside matters as much as the upside.")
img(s,"assets/03_asymmetry.png",5.9,2.2,6.9,4.3)
bullets(s,0.57,2.5,5.1,3.4,[
 "Fix Operations: +₱0.2M → +₱3.5M — the recurring floor stays positive.",
 "Cebu hub (owned): −₱6.9M → +₱1.0M — loses in 2 of 3.",
 "Mall kiosk: −₱2.5M → +₱3.0M — a coin-flip on foot traffic (Assumption F).",
 "Decision rule: protect the downside first; earn the upside through measured pilots."],sz=13,gap=12)
source(s,"Model on Exhibit C + Assumptions C/D/E/F · full math in Appendix")
footer(s); notes(s,"We handicapped the hub generously (50% margin) and discounted our own plan (30%) — Fix-Ops still wins.")

# ============ S7 · C — STRATEGY ============
s=slide(); title2(s,[("FIX THE ENGINE FIRST, ",TITLE),("THEN EXPAND",LIME)],sect="C")
voice(s,"One primary move, one supporting channel test, one gated pilot — each step must earn the next.")
def phase(x,tag,tagc,head,body):
    rect(s,x,2.45,3.9,2.6,SLATE,line=tagc,lw=2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,x+0.25,2.65,3.5,0.4,[(tag,12,True,tagc)])
    text(s,x+0.25,3.1,3.5,0.5,[(head,16.5,True,SOFT)])
    text(s,x+0.25,3.8,3.5,1.2,[(body,11.5,False,INK)])
phase(0.57,"PRIMARY · MONTHS 1–6",LIME,"Fix operations","Build stock visibility + earlier warnings. Test a small webstore and measure real orders.")
phase(4.72,"GATED · MONTH 7+",GOLD,"Cebu 3PL pilot","Only if all partner costs fit below ~₱56/order — then prove it in an 8-week pilot.")
phase(8.87,"YEAR 2 · DEFERRED",MUTE,"Kiosk / owned hub","Evaluated later from real operating and Cebu-pilot evidence; not funded in Year 1.")
rich(s,0.57,5.5,12.2,0.7,[("The principle:  ",14.5,True,LIME),("earn the right to expand. Year 1 builds "
    "operating proof that makes Year-2 growth safer.",14.5,True,SOFT)],align=PP_ALIGN.CENTER)
footer(s); notes(s,"Resolves the 'do everything' trap the case penalizes. One real investment + a near-free hedge + a gated pilot.")

# ============ S8 · C — THE SYSTEM ============
s=slide(); title2(s,[("AN EARLIER WARNING, ",TITLE),("WHILE THERE'S TIME",LIME)],sect="C")
voice(s,"Ramen can have 20 days of stock left and still be late to reorder — because its supplier needs 30 days.")
steps=[("1 · UPDATE","Bring daily sales, current stock, incoming orders, and supplier lead time into one sheet.",LIME),
       ("2 · COMPARE","Ask: will available and incoming stock last until the next delivery arrives?",SKY),
       ("3 · WARN","Flag the products likely to run short and show the projected gap.",GOLD),
       ("4 · DECIDE","Founder checks cash, supplier availability, and minimum order size before approving.",ORANGE)]
x=0.57
for tag,body,c in steps:
    rect(s,x,2.5,2.92,2.45,SLATE,line=BORDER,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s,x,2.5,2.92,0.07,c); text(s,x+0.2,2.72,2.6,0.4,[(tag,12.5,True,c)])
    text(s,x+0.2,3.2,2.6,1.6,[(body,11,False,INK)]); x+=3.06
rect(s,0.57,5.3,12.2,0.72,SLATE2,line=GREEN,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rich(s,0.8,5.4,12.0,0.55,[("5 · LEARN   ",13,True,GREEN),("Compare alerts with actual results weekly. Start with Excel + Power Query + barcode checks; "
    "add AI trend signals only after the basic data is reliable.",12.2,False,SOFT)],anchor=MSO_ANCHOR.MIDDLE)
source(s,"Exhibit D (current process) · Exhibit C (lead times)")
footer(s); notes(s,"Highest-weighted criterion (×5): real workflow, real tools, human-in-the-loop.")

# ============ S9 · C — PROTOTYPE ============
s=slide(); title2(s,[("WE TESTED THE ",TITLE),("DECISION PROCESS",LIME)],sect="C")
voice(s,"A working prototype shows how an earlier warning changes when the evidence changes.")
img(s,"screenshot_3_dashboard.png",2.1,2.35,9.1,3.7,frame=True,
    caption="Prototype — case baselines + clearly labeled illustrative operating data")
text(s,0.57,6.12,12.2,0.55,[("Example: ramen shows ~20 days of stock cover, but its supplier needs 30 days (Exhibit C) — so the warning appears before the shelf is empty. "
    "This is not PopCart's live inventory and it does not place orders automatically.",11.2,False,MUTE)],align=PP_ALIGN.CENTER)
footer(s); notes(s,"The credibility spike — the slide no other team can fake.")

# ============ S10 · C — WHAT YEAR 1 SOLVES / DEFERS (Q4) ============
s=slide(); title2(s,[("WHAT YEAR 1 SOLVES — ",TITLE),("AND DEFERS",LIME)],sect="C")
voice(s,"A credible strategy states what the primary move improves — and what still needs a separate pilot.")
cols=[("PROBLEM",0.57,3.0),("BASELINE (SOURCE)",3.7,3.5),("AFTER YEAR 1",7.4,2.5),("HOW",10.0,2.8)]
y=2.45
rect(s,0.57,y,12.2,0.5,SLATE2)
for h,x,w in cols: text(s,x+0.12,y+0.07,w,0.4,[(h,11.5,True,LIME)])
y+=0.5
data=[("Stockouts","11 days/mo (Exhibit C)","~6 base; ~4–8 range","earlier warnings + 25–60% reduction"),
      ("Inventory visibility","2×/day manual (Exhibit D)","every 2h after pilot","scheduled refresh + barcode checks"),
      ("Channel risk","90% TikTok (Exhibit D)","second channel tested","webstore sales must be earned, not assumed"),
      ("Shipping / delay / damage","₱155, 5–7d, 4.2% CV (Exhibit B)","not promised by primary move","only after a successful 3PL pilot")]
for i,(a,b,c,d) in enumerate(data):
    fc=SLATE if i%2 else BG
    rect(s,0.57,y,12.2,0.78,fc)
    text(s,0.69,y+0.04,3.0,0.7,[(a,12.5,True,SOFT)],anchor=MSO_ANCHOR.MIDDLE)
    text(s,3.82,y+0.04,3.5,0.7,[(b,11.5,False,INK)],anchor=MSO_ANCHOR.MIDDLE)
    text(s,7.52,y+0.04,2.5,0.7,[(c,11.5,True,GREEN)],anchor=MSO_ANCHOR.MIDDLE)
    text(s,10.12,y+0.04,2.7,0.7,[(d,11,False,INK)],anchor=MSO_ANCHOR.MIDDLE)
    y+=0.78
source(s,"Exhibit B · Exhibit C · Exhibit D · Assumption D",y=6.7)
footer(s); notes(s,"Directly answers Q4 and proves we read every exhibit, not just the convenient ones.")

# ============ S11 · C — WHERE & WHEN ============
s=slide(); title2(s,[("TEST CEBU BEFORE ",TITLE),("OWNING THE HUB",LIME)],sect="C")
voice(s,"A 3PL is an outside warehouse that stores and ships selected products for a per-order fee.")
bullets(s,0.57,2.45,5.6,3.6,[
 "Owned hub: ₱7.2–10.8M/yr operating cost whether volume arrives or not (Assumption C).",
 "3PL pilot: smaller commitment, per-order charges, selected products only.",
 "Cebu before Davao: higher volume (26% vs 17% — Exhibit B / Assump. B).",
 "Gate: storage + transfer + pick-pack + returns + minimums must fit below ~₱56/order.",
 "Proof: 8-week pilot; measure total cost/order, delivery time, damage, and stock availability."],sz=12.2,gap=9)
img(s,"assets/04_grab_vs_car.png",6.5,2.45,6.3,4.0)
source(s,"Exhibit B · Exhibit D · Assumptions B/C/D")
footer(s); notes(s,"Q2 fully answered: where (Cebu via 3PL) and when (only when the unit economics earn it).")

# ============ S12 · D — HONEST MONEY ============
s=slide(); title2(s,[("THE ",TITLE),("HONEST MONEY",LIME)],sect="D")
voice(s,"We could have claimed the full ₱7.35M. Here's the number we'd actually keep.")
img(s,"assets/02_lost_profit_by_sku.png",0.57,2.4,6.1,3.6)
img(s,"assets/06_contribution_waterfall.png",6.9,2.4,5.9,3.6)
text(s,0.57,5.98,12.2,0.78,[("Base math: ₱14.70M lost sales × 42.5% recovery = ₱6.25M recovered revenue. "
    "At 30–40% money left after product, platform, and shipping costs: ₱1.87–2.50M. "
    "Less ₱0.90M annual system cost = ₱0.97–1.60M recurring benefit; one-time setup = ₱1.40M.",
    10.6,False,MUTE)],align=PP_ALIGN.CENTER)
source(s,"Exhibit A · Exhibit C · Assumption D",y=6.72)
footer(s); notes(s,"Honesty is scored. Discounting our own upside is what makes the number believable.")

# ============ S13 · D — PAYBACK ============
s=slide(); title2(s,[("PAYBACK, ",TITLE),("WITHOUT PRETENDING",LIME)],sect="D")
voice(s,"Expected payback is ~13–20 months at 30–40% take-home margin — slower if fewer sales are recovered.")
img(s,"assets/07_payback_sensitivity.png",0.57,2.45,6.3,3.6)
img(s,"assets/08_investment_tree.png",6.95,2.45,5.9,3.6)
text(s,0.57,6.1,12.2,0.55,[("Year 1 is an investment year (−₱0.67M cash); net-positive thereafter. The case allows a payback "
    "beyond 12 months if stated honestly. The hub never pays back in its base case.",11.5,False,MUTE)],align=PP_ALIGN.CENTER)
source(s,"Model on Exhibit C + Assumptions C/D",y=6.72)
footer(s); notes(s,"Pre-empts the toughest finance question before a judge can raise it.")

# ============ S14 · E — ROADMAP ============
s=slide(); title2(s,[("BUILD THE HABIT, ",TITLE),("THEN TEST EXPANSION",LIME)],sect="E")
voice(s,"Each phase produces evidence for the next decision — with a clear stop point if the numbers fail.")
rect(s,0.9,4.0,11.55,0.04,BORDER)
nodes=[(1.25,"M1","BASELINE","Count stock; clean SKU + supplier list; confirm lead times.",LIME),
       (3.35,"M2–3","PILOT","5 fast movers; daily warnings; log founder decisions.",SKY),
       (5.45,"M4–6","STABILIZE","2h updates; target 6–8 stockout days; webstore test.",GREEN),
       (7.55,"M7 · GATE","QUOTE","All-in cost <~₱56/order; cash available for split stock.",ORANGE),
       (9.65,"M8–9","IF GO: PILOT","8 weeks; track cost, speed, damage, and availability.",GOLD),
       (11.75,"M10–12","DECIDE","SOPs + Year-2 go/no-go on expansion and kiosks.",GREEN)]
for x,m,head,body,c in nodes:
    if "GATE" in m: rect(s,x-0.16,3.84,0.32,0.32,ORANGE,shape=MSO_SHAPE.DIAMOND)
    else: rect(s,x-0.12,3.88,0.24,0.24,c,shape=MSO_SHAPE.OVAL)
    text(s,x-0.92,3.22,1.84,0.42,[(m,11.5,True,c)],align=PP_ALIGN.CENTER)
    text(s,x-0.92,4.32,1.84,0.42,[(head,10.8,True,SOFT)],align=PP_ALIGN.CENTER)
    text(s,x-0.96,4.72,1.92,1.15,[(body,9.2,False,INK)],align=PP_ALIGN.CENTER)
text(s,0.75,6.1,12.0,0.5,[("People: warehouse manager + two admins run the daily process; founders approve purchase orders and the Cebu go/no-go.",
    11.2,True,SOFT)],align=PP_ALIGN.CENTER)
footer(s); notes(s,"The Month-7 diamond is the proof of discipline the case rewards.")

# ============ S15 · E — RISKS, IMPACT & CLOSE ============
s=slide(); title2(s,[("STOP THE LEAK. ",TITLE),("EARN THE GROWTH.",LIME)],sect="E")
voice(s,"We tried to break our own plan — here's where it held, and where it lands.")
text(s,0.57,2.35,6.0,0.4,[("RISKS → MITIGATIONS",12,True,LIME)])
bullets(s,0.57,2.8,6.0,3.2,[
 "Staff adoption → familiar Excel/Power Query workflow + training.",
 "Channel shock → webstore test by Month 6; measure orders before claiming diversification.",
 "Recovery below model → track stockout days + recovered sales monthly; adjust or stop.",
 "Hub under-use → deferred behind a unit-economics gate.",
 "Long leads, ramen 30d (Exhibit C) → pre-order + 2nd supplier."],sz=12,gap=9)
text(s,6.9,2.35,6.0,0.4,[("EXPECTED IMPACT",12,True,LIME)])
statcard(s,6.9,2.8,1.85,1.95,"~6",GREEN,"stockout days","from 11 (Exhibit C)")
statcard(s,8.95,2.8,1.85,1.95,"2h",SKY,"update cycle","from 2×/day (Exh. D)")
statcard(s,11.0,2.8,1.75,1.95,"₱1M",LIME,"recurring / yr","base ₱0.97M; up to ₱1.6M")
rect(s,6.9,5.0,5.85,1.15,SLATE,line=LIME,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,7.1,5.07,5.5,1.0,[("“Success isn't found in the newest hub —",15,True,SOFT),
    ("it's in the precision of the shelf.”",15,True,LIME)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
footer(s); notes(s,"Red-team + impact + payoff line that ties back to the cover. The engine, fixed.")

# ============ APPENDIX A ============
s=slide(); title2(s,[("APPENDIX A — ",TITLE),("AI DISCLOSURE",LIME)])
voice(s,"AI was our analytical partner — and we checked every number it produced against the exhibits.")
rows=[("AI assistants ([Claude / Gemini / ChatGPT])","Structured the exhibits; ranked the money leaks; drafted copy; red-teamed the strategy"),
 ("Python (written with AI help)","popcart_model.py — scenario financials & payback;  build_prototype.py — the reorder prototype"),
 ("Google Trends / public web","Market context on TikTok Shop PH & snack e-commerce"),
 ("[Canva / PowerPoint]","Slide design, charts, and this deck")]
y=2.5; rect(s,0.57,y,4.8,0.5,SLATE2); text(s,0.75,y+0.07,4.6,0.4,[("TOOL",12,True,LIME)])
rect(s,5.6,y,7.17,0.5,SLATE2); text(s,5.78,y+0.07,6.9,0.4,[("WHERE WE USED IT",12,True,LIME)]); y+=0.5
for i,(a,b) in enumerate(rows):
    fc=SLATE if i%2 else BG; rect(s,0.57,y,4.8,0.72,fc); rect(s,5.6,y,7.17,0.72,fc)
    text(s,0.75,y+0.04,4.6,0.65,[(a,11,True,SOFT)],anchor=MSO_ANCHOR.MIDDLE)
    text(s,5.78,y+0.04,6.9,0.65,[(b,11,False,INK)],anchor=MSO_ANCHOR.MIDDLE); y+=0.72
text(s,0.57,y+0.15,12.2,0.5,[('"AI was an analytical and drafting partner. Every figure was verified against '
    'the case exhibits. No output was presented as fact without human checking."',11.5,False,MUTE)],italic=True)
footer(s); notes(s,"Fill in the exact tools your team used before submitting.")

# ============ APPENDIX B ============
s=slide(); title2(s,[("APPENDIX B — ",TITLE),("THE WORKING MODEL",LIME)])
voice(s,"Every cell is a real formula — the README tab documents the math, so it's fully auditable.")
img(s,"screenshot_2_control_panel.png",0.57,2.4,12.2,3.3,frame=True,
    caption="Reorder Control Panel — forecast, safety stock, reorder point, status, suggested qty (live formulas)")
text(s,0.57,6.0,12.2,0.9,[("Forecast/day = weighted 7/14/28-day × viral multiplier   ·   "
    "Safety stock = 1.65 × volatility × √lead-time (~95% service)   ·   Reorder point = forecast × lead-time + safety stock",
    11,False,MUTE)],align=PP_ALIGN.CENTER)
footer(s); notes(s,"Upload PopCart_Reorder_System.xlsx to Google Sheets to demo live if you reach the semifinal.")

# ============ APPENDIX C ============
s=slide(); title2(s,[("APPENDIX C — ",TITLE),("ASSUMPTIONS & SOURCES",LIME)])
voice(s,"We'd rather show our working than hide it — here is every assumption behind the numbers.")
bullets(s,0.57,2.5,12.0,4.0,[
 "AOV ₱650; 185,000 orders/yr (Assumption A/B). Blended gross margin 50% (Exhibit A).",
 "Contribution margin ~30% on recovered sales — nets out platform fees + shipping (the honest value, not 50% gross).",
 "Stockout recovery 25–60% (Assumption D); base case 42.5%.",
 "Damage loss/order ~₱480 (COGS + wasted shipping).",
 "Hub growth credited generously at 50% — we handicap the option we reject; Fix-Ops still wins.",
 "Expansion costs from Assumption C; 3-month build before benefits flow; tools at MSME low-end.",
 "Prototype daily-sales & TikTok index are illustrative samples — the logic is the deliverable."],sz=12.5,gap=10)
footer(s); notes(s,"State assumptions clearly and disclose AI — both explicitly rewarded by the rubric.")

prs.save("PopCart_PH_LAMBO2026_DRAFT.pptx")
print("Saved PopCart_PH_LAMBO2026_DRAFT.pptx —", len(prs.slides._sldIdLst), "slides")
